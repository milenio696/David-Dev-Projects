import os
from enum import Enum

import azure.identity
import openai
import rich
from dotenv import load_dotenv
from pydantic import BaseModel, Field

# Librerías de lectura para Office 365
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


# ==============================
# CONFIGURACIÓN DEL CLIENTE

load_dotenv(override=True)
API_HOST = os.getenv("API_HOST", "github")

if API_HOST == "azure":
    token_provider = azure.identity.get_bearer_token_provider(
        azure.identity.DefaultAzureCredential(),
        "https://cognitiveservices.azure.com/.default"
    )
    client = openai.OpenAI(
        base_url=os.environ["AZURE_OPENAI_ENDPOINT"],
        api_key=token_provider,
    )
    MODEL_NAME = os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT"]

elif API_HOST == "ollama":
    client = openai.OpenAI(
        base_url=os.environ["OLLAMA_ENDPOINT"],
        api_key="nokeyneeded",
    )
    MODEL_NAME = os.environ["OLLAMA_MODEL"]

elif API_HOST == "github":
    client = openai.OpenAI(
        base_url="https://models.github.ai/inference",
        api_key=os.environ["GITHUB_TOKEN"],
    )
    MODEL_NAME = os.getenv("GITHUB_MODEL", "openai/gpt-4o")

else:
    client = openai.OpenAI(api_key=os.environ["OPENAI_KEY"])
    MODEL_NAME = os.environ["OPENAI_MODEL"]

rich.print(f"[bold green]Usando modelo:[/bold green] {MODEL_NAME} ({API_HOST})")


# ==============================
# MODELOS ESTRUCTURADOS

class CalendarEvent(BaseModel):
    name: str
    date: str = Field(..., description="Fecha del evento (YYYY-MM-DD)") # Se setea el formato deseado
    participants: list[str] 


class DeliveryQuery(BaseModel):
    order_id: str


class DocumentSummary(BaseModel):
    filename: str
    summary: str
    keywords: list[str]


# ==============================
# FUNCIONES AUXILIARES 

def read_office_document(filepath: str) -> str:
    """
    Lee el contenido textual de un archivo de Office (.docx, .xlsx o .pptx)
    y lo devuelve como una cadena de texto unificada.

    Soporta tres tipos de archivo:
        - Word (.docx): concatena el texto de todos los párrafos.
        - Excel (.xlsx): concatena el contenido de cada celda, hoja por hoja.
        - PowerPoint (.pptx): concatena el texto de todas las diapositivas y formas con texto.
    """
     # Obtiene la extensión del archivo y la convierte a minúsculas
    ext = os.path.splitext(filepath)[1].lower()  # '.docx', '.xlsx', '.pptx'

    if ext == ".docx":
        # Abre el documento Word usando python-docx
        doc = Document(filepath)

        # Une todos los párrafos en una sola cadena, separando cada párrafo con un salto de línea
        text_content = "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".xlsx":
        # Abre el libro Excel usando openpyxl
        wb = load_workbook(filepath, data_only=True)
        text_content = ""  # Inicializa la variable de texto

        # Recorre todas las hojas del libro
        for sheet in wb:
            # Recorre cada fila y concatena los valores de las celdas en una línea
            for row in sheet.iter_rows(values_only=True):
                text_content += " ".join(str(cell) for cell in row if cell) + "\n"

    elif ext == ".pptx":
        # Abre la presentación PowerPoint usando python-pptx
        prs = Presentation(filepath)
        text_content = ""  # Inicializa la variable de texto

        # Recorro las diapositivas
        for slide in prs.slides:
            # Recorre cada forma dentro de la diapositiva
            for shape in slide.shapes:
                # Solo toma las formas que tengan texto
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"

    else:
        # Si el archivo no es de los permitidos, lanza un error
        raise ValueError(f"Tipo de archivo no soportado: {ext}")

    # Retorna el texto, eliminando espacios en blanco al inicio y final
    return text_content.strip()



# ==============================
# AUTOMATIZACIÓN DE TAREAS

tareas = [
    {
        "nombre": "Eventos",
        "prompt": "David y Leonardo van a una feria de programación el 16 de abril.",
        "schema": CalendarEvent,
        "instruccion": "Extrae la información del evento. Si no hay año, imagina que es 2025."
    },
    {
        "nombre": "Consulta de envio",
        "prompt": "¿Cuándo llegará mi pedido #84975903?",
        "schema": DeliveryQuery,
        "instruccion": "Sos un asistente de atención al cliente. Extrae el número de pedido y responde brevemente."
    },
    {
        "nombre": "Lectura de documentos",
        "prompt": None,  # Se completará mas adelante
        "schema": DocumentSummary,
        "instruccion": "Resume el contenido del documento dado y da un resumen pequeño de lo que fue el documento."
    },
]

# Ruta de prueba del documento de Office. No se encuentra en el repositorio
DOC_PATH = "../sample_doc.docx"

# Si existe, agrega el contenido leído al prompt
if os.path.exists(DOC_PATH):
    doc_text = read_office_document(DOC_PATH)
    tareas[2]["prompt"] = f"Analiza el siguiente documento:\n\n{doc_text[:4000]}"  # Se estalece el prompt junto con su limite de texto
else:
    tareas.pop(2)  # Si no hay documento, omite la tarea


# ==============================
# EJECUCIÓN AUTOMÁTICA

for tarea in tareas:
    # Imprime en consola el nombre de la tarea que se está ejecutando, con formato azul y negrita
    rich.print(f"\n[bold blue]--- Ejecutando: {tarea['nombre']} ---[/bold blue]")

    # Llama a la función 'run_structured_task' para enviar la instrucción al modelo y tener la salida estructurada del schema
    resultado = run_structured_task(
        prompt=tarea["prompt"],              # Texto o instrucción del usuario para esta tarea
        schema=tarea["schema"],              # Esquema Pydantic que valida la estructura de la respuesta
        system_instruction=tarea["instruccion"]  # Instrucción del sistema para guiar al modelo
    )

    # Si el modelo devolvió un resultado válido (no se negó a responder)
    if resultado:
        # Muestra el resultado estructurado en la consola
        rich.print(resultado)